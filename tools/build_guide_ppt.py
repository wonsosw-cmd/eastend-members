# -*- coding: utf-8 -*-
"""EASTEND 오프라인 멤버십 — 매장 운영 가이드 PPT (v2.4: 웰컴 3,000P + 익일 적립)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
AMBER = RGBColor(0xE9, 0xA1, 0x3B)
BG = RGBColor(0xF6, 0xF7, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x22, 0x29, 0x3A)
MUTED = RGBColor(0x6B, 0x74, 0x86)
LINE = RGBColor(0xE4, 0xE8, 0xEE)
TEAL_T = RGBColor(0xE2, 0xF3, 0xF1)
AMBER_T = RGBColor(0xFC, 0xF0, 0xDE)
NAVY_T = RGBColor(0xE8, 0xEC, 0xF5)
RED = RGBColor(0xD6, 0x45, 0x45)

FONT = "맑은 고딕"
QR_SAMPLE = r"G:\공유 드라이브\영업마케팅실\★. 데일리 매출 보고\★오프라인 회원관리\QR_포스터용_PNG\C005_시티브리즈_플래그십.png"
OUT = r"G:\공유 드라이브\영업마케팅실\★. 데일리 매출 보고\★오프라인 회원관리\EASTEND_멤버십_매장가이드_v1.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=None, line=None, radius=0.08):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def circle(s, x, y, d, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.15, space_after=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, {})]
        para = [(it, {}) if isinstance(it, str) else it for it in para]
        for t, st in para:
            r = p.add_run()
            r.text = t
            r.font.name = FONT
            r.font.size = Pt(st.get("size", size))
            r.font.bold = st.get("bold", bold)
            r.font.color.rgb = st.get("color", color)
            rPr = r._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {})
                rPr.append(ea)
            ea.set('typeface', FONT)
    return tb


def num_badge(s, x, y, n, color=TEAL, d=0.42):
    circle(s, x, y, d, color)
    text(s, x, y + 0.015, d, d - 0.03, str(n), size=17, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def title_bar(s, step, title, sub=""):
    text(s, 0.65, 0.42, 1.9, 0.35, step, size=13, color=TEAL, bold=True)
    text(s, 0.65, 0.72, 10.5, 0.75, title, size=30, color=NAVY, bold=True)
    if sub:
        text(s, 0.65, 1.38, 12.0, 0.4, sub, size=13.5, color=MUTED)


# ───────────────────────── 1. 표지
s = slide(NAVY)
text(s, 0.9, 1.55, 6.0, 0.4, "CITYBREEZE  ×  ARTID", size=15, color=AMBER, bold=True)
text(s, 0.9, 2.05, 11.5, 1.9,
     [[("EASTEND 오프라인 멤버십", {})], [("매장 운영 가이드", {})]],
     size=44, color=WHITE, bold=True, line_spacing=1.1)
text(s, 0.9, 4.15, 11.0, 0.9,
     [["QR 가입  ·  영수증 적립  ·  마일리지 에누리"],
      [("전 매장 공통 절차 안내  |  2026. 07", {"color": RGBColor(0xAF, 0xB9, 0xCE), "size": 14})]],
     size=17, color=WHITE)
stats = [("신규가입 혜택", "3,000P 즉시 지급"), ("적립률", "구매액의 1% 자동"), ("사용 (1P=1원)", "5,000P부터 1,000P 단위")]
for i, (k, v) in enumerate(stats):
    x = 0.9 + i * 4.0
    box(s, x, 5.45, 3.6, 1.25, fill=RGBColor(0x25, 0x36, 0x5C))
    text(s, x + 0.3, 5.68, 3.0, 0.3, k, size=12.5, color=RGBColor(0xAF, 0xB9, 0xCE), bold=True)
    text(s, x + 0.3, 6.02, 3.2, 0.5, v, size=19, color=WHITE, bold=True)

# ───────────────────────── 2. 한눈에 보기
s = slide()
title_bar(s, "OVERVIEW", "멤버십, 이렇게 돌아갑니다",
          "고객은 QR로 가입하고, 구매하면 자동으로 적립되고, 매장에서 포인트만큼 에누리 받습니다.")
flow = [
    ("1", "QR 스캔 → 가입", "매장 QR 포스터를 스캔해\n정보 입력·동의하면 가입 완료\n신규가입 3,000P 즉시 지급", TEAL),
    ("2", "구매 시 자동 적립", "구매하면 결제금액의 1%가\n자동으로 적립\n(회원가입 익일부터)", AMBER),
    ("3", "1년간 사용 가능", "포인트는 적립일로부터\n1년간 유효하며\n미사용분은 자동 소멸", NAVY),
    ("4", "매장에서 에누리", "결제 시 전화번호 조회 후\n보유 포인트만큼 차감\n(1P = 1원)", TEAL),
]
for i, (n, t, d, c) in enumerate(flow):
    x = 0.65 + i * 3.13
    box(s, x, 2.0, 2.85, 3.3, fill=WHITE, line=LINE)
    num_badge(s, x + 0.28, 2.3, n, color=c)
    text(s, x + 0.28, 2.95, 2.35, 0.75, t, size=16.5, color=NAVY, bold=True)
    text(s, x + 0.28, 3.72, 2.35, 1.45, d, size=12, color=MUTED, line_spacing=1.3)
    if i < 3:
        text(s, x + 2.83, 3.35, 0.35, 0.5, "→", size=20, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
box(s, 0.65, 5.65, 12.05, 1.15, fill=NAVY_T)
text(s, 1.0, 5.92, 11.4, 0.7,
     [[("매장에서 하는 일은 딱 2가지  —  ", {"bold": True, "color": NAVY}),
       ("① QR 가입 안내   ② 포인트 차감(에누리)  — 적립은 자동!", {"color": TEXT})]],
     size=15)

# ───────────────────────── 3. STEP 1 고객 가입
s = slide()
title_bar(s, "STEP 1", "고객 가입 안내 — QR만 보여주세요",
          "가입은 고객 본인 휴대폰에서 진행됩니다. 매장은 QR 포스터만 안내하면 됩니다.")
steps = [
    ("QR 스캔", "계산대·피팅룸 옆에 비치된 매장 전용 QR을 고객 휴대폰 카메라로 스캔"),
    ("정보 입력", "이름·휴대전화번호·생년월일(필수), 성별(선택) 입력"),
    ("약관 동의", "필수/선택 항목을 구분해 동의 — ‘전체 동의’는 편의 기능(개별 해제 가능)"),
    ("가입 완료", "신규가입 3,000P 즉시 지급 — 구매 적립(1%)은 가입 다음 날부터 자동"),
]
for i, (t, d) in enumerate(steps):
    y = 2.0 + i * 1.15
    num_badge(s, 0.75, y + 0.08, i + 1, color=TEAL)
    text(s, 1.4, y, 6.3, 0.35, t, size=16, color=NAVY, bold=True)
    text(s, 1.4, y + 0.38, 6.5, 0.65, d, size=12.5, color=MUTED, line_spacing=1.25)
box(s, 8.6, 1.95, 4.1, 4.75, fill=WHITE, line=LINE)
try:
    s.shapes.add_picture(QR_SAMPLE, Inches(9.55), Inches(2.35), Inches(2.2), Inches(2.2))
except Exception:
    pass
text(s, 8.85, 4.7, 3.6, 0.4, "매장별 전용 QR 포스터", size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
text(s, 8.85, 5.12, 3.6, 1.3,
     [["매장마다 QR이 다릅니다 (매장 실적 집계)."],
      ["포스터 인쇄본은 본사에서 배포하며,"],
      ["훼손 시 본사로 재요청해주세요."]],
     size=11.5, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.25)
box(s, 0.65, 6.35, 7.6, 0.75, fill=AMBER_T)
text(s, 0.95, 6.52, 7.1, 0.45,
     [[("멘트 예시  ", {"bold": True, "color": NAVY}),
       ("“지금 QR 찍고 가입하시면 3,000P를 바로 드려요!”", {"color": TEXT})]], size=13)

# ───────────────────────── 4. STEP 2 적립 안내 (자동)
s = slide()
title_bar(s, "STEP 2", "적립은 자동입니다 — 매장이 할 일 없음",
          "구매 시 판매 프로그램에서 결제금액의 1%가 자동 적립됩니다. 별도 등록·입력이 필요 없습니다.")
box(s, 0.65, 1.95, 5.95, 3.0, fill=WHITE, line=LINE)
text(s, 1.0, 2.2, 5.2, 0.4, "🎁 신규가입 웰컴", size=16, color=TEAL, bold=True)
text(s, 1.0, 2.72, 5.3, 2.0,
     [["QR 가입이 완료되는 순간"],
      [("3,000P가 즉시 자동 지급", {"bold": True, "color": NAVY})], [""],
      ["매장은 가입 안내만 하면 되고,"],
      ["지급은 시스템이 자동 처리합니다."]],
     size=13, color=TEXT, line_spacing=1.3)
box(s, 6.75, 1.95, 5.95, 3.0, fill=WHITE, line=LINE)
text(s, 7.1, 2.2, 5.2, 0.4, "🛍 구매 적립", size=16, color=AMBER, bold=True)
text(s, 7.1, 2.72, 5.3, 2.0,
     [["구매(결제) 시"],
      [("결제금액의 1%가 자동 적립", {"bold": True, "color": NAVY})], [""],
      [("단, 회원가입 익일부터 적용", {"bold": True, "color": RED})],
      ["(가입 당일 구매분은 적립되지 않습니다)"]],
     size=13, color=TEXT, line_spacing=1.3)
box(s, 0.65, 5.2, 12.05, 1.75, fill=NAVY_T)
text(s, 1.0, 5.4, 11.5, 0.35, "적립 규칙", size=14, color=NAVY, bold=True)
text(s, 1.0, 5.78, 11.5, 1.1,
     [[("·  적립액 = 결제금액의 1% (원단위 절사, 구매 시 자동)      ·  별도 등록·영수증 입력 불필요", {})],
      [("·  적립은 ", {}), ("회원가입 익일부터", {"bold": True, "color": RED}),
       ("  가능 (가입 당일 적립 불가 — 대신 웰컴 3,000P 즉시 지급)", {})],
      [("·  포인트는 ", {}), ("적립일로부터 1년", {"bold": True, "color": RED}),
       ("  경과 시 자동 소멸  ·  반품 시 적립 취소는 본사에서 처리", {})]],
     size=13, color=TEXT, line_spacing=1.3)

# ───────────────────────── 5. STEP 3 매장용 페이지 설정
s = slide()
title_bar(s, "STEP 3", "매장용 페이지 최초 설정 (1회)",
          "매장 태블릿·공용폰 크롬에서 한 번만 설정하면 이후 자동 로그인됩니다.")
box(s, 0.65, 1.95, 7.55, 4.4, fill=WHITE, line=LINE)
setup = [
    ("접속", "크롬에서 아래 주소 입력 (즐겨찾기/홈화면 추가 권장)"),
    ("매장용 토큰 입력", "아래 토큰을 정확히 입력"),
    ("우리 매장 선택", "브랜드·매장명 목록에서 우리 매장 선택"),
    ("시작하기", "이 기기에 저장됨 — 다음부터 바로 조회 화면"),
]
for i, (t, d) in enumerate(setup):
    y = 2.25 + i * 1.0
    num_badge(s, 1.0, y + 0.04, i + 1, color=NAVY)
    text(s, 1.65, y, 6.3, 0.35, t, size=15.5, color=NAVY, bold=True)
    text(s, 1.65, y + 0.37, 6.4, 0.55, d, size=12.5, color=MUTED)
box(s, 8.6, 1.95, 4.1, 4.4, fill=NAVY)
text(s, 8.95, 2.3, 3.4, 0.35, "매장용 페이지 주소", size=12.5, color=RGBColor(0xAF, 0xB9, 0xCE), bold=True)
text(s, 8.95, 2.68, 3.5, 0.9, "wonsosw-cmd.github.io\n/eastend-members/staff.html",
     size=14.5, color=WHITE, bold=True, line_spacing=1.2)
text(s, 8.95, 3.85, 3.4, 0.35, "매장용 토큰", size=12.5, color=RGBColor(0xAF, 0xB9, 0xCE), bold=True)
box(s, 8.95, 4.22, 3.4, 0.62, fill=RGBColor(0x25, 0x36, 0x5C))
text(s, 8.95, 4.32, 3.4, 0.45, "eastlove", size=18, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
text(s, 8.95, 5.05, 3.45, 1.1,
     [["토큰은 매장 내부용입니다."], ["고객·외부에 노출되지 않도록"], ["주의해주세요."]],
     size=11.5, color=RGBColor(0xAF, 0xB9, 0xCE), line_spacing=1.3)
box(s, 0.65, 6.55, 12.05, 0.6, fill=AMBER_T)
text(s, 1.0, 6.68, 11.4, 0.4,
     [[("TIP  ", {"bold": True, "color": AMBER}),
       ("크롬 메뉴 → ‘홈 화면에 추가’를 해두면 앱처럼 바로 열 수 있습니다.", {"color": TEXT})]], size=12.5)

# ───────────────────────── 6. STEP 4 조회·적립·차감
s = slide()
title_bar(s, "STEP 4", "일상 운영 — 조회 · 차감",
          "모든 작업은 ‘전화번호 조회 → 고객 확인’ 후에 진행합니다. 적립은 자동입니다.")
box(s, 0.65, 1.95, 5.95, 4.4, fill=WHITE, line=LINE)
text(s, 1.0, 2.2, 5.2, 0.4, "① 조회 & 본인 확인", size=15.5, color=NAVY, bold=True)
text(s, 1.0, 2.72, 5.3, 3.4,
     [["고객 전화번호 입력 → 조회"],
      [""],
      ["화면에 이름이 홍*동 처럼 일부만 표시됩니다."],
      [""],
      [("“홍길동 고객님 맞으세요?”", {"bold": True, "color": NAVY})],
      ["확인 후 진행해주세요."],
      [""],
      [("잔액(사용 가능 마일리지)이 함께 표시됩니다.", {"color": MUTED, "size": 11.5})]],
     size=12.5, color=TEXT, line_spacing=1.3)
box(s, 6.75, 1.95, 5.95, 4.4, fill=AMBER_T)
text(s, 7.1, 2.2, 5.2, 0.4, "② 마일리지 차감 (에누리)", size=15.5, color=RGBColor(0xB0, 0x72, 0x1A), bold=True)
text(s, 7.1, 2.72, 5.3, 3.4,
     [["‘마일리지 차감’에서 차감할 포인트 입력 (1P = 1원)"],
      [""],
      [("→ 5,000P부터, 1,000P 단위", {"bold": True, "color": NAVY})],
      [("→ 결제 금액에서 에누리 적용", {"bold": True, "color": NAVY})],
      [""],
      [("잔액 초과·규칙 위반 시 차감되지 않고", {"color": MUTED, "size": 11.5})],
      [("안내 문구가 표시됩니다.", {"color": MUTED, "size": 11.5})],
      [""],
      [("적립은 자동이므로 매장 입력이 없습니다.", {"color": MUTED, "size": 11.5})]],
     size=12.5, color=TEXT, line_spacing=1.3)
box(s, 0.65, 6.55, 12.05, 0.6, fill=NAVY_T)
text(s, 1.0, 6.68, 11.4, 0.4,
     [[("주의  ", {"bold": True, "color": RED}),
       ("차감은 되돌리기 어렵습니다. 금액을 고객과 함께 확인한 뒤 진행해주세요.", {"color": TEXT})]], size=12.5)

# ───────────────────────── 7. FAQ
s = slide()
title_bar(s, "FAQ", "자주 묻는 질문 & 꼭 지켜주세요")
faqs = [
    ("포인트는 언제 쌓이나요?", "구매 시 결제금액의 1%가 자동 적립됩니다 (회원가입 익일부터). 매장에서 별도 등록할 것이 없습니다."),
    ("미가입 고객이에요", "QR 가입을 먼저 안내해주세요. 가입 즉시 웰컴 3,000P가 자동 지급됩니다."),
    ("가입 당일 구매도 적립되나요?", "아니요. 적립은 회원가입 다음 날 구매분부터 적용됩니다. 대신 웰컴 3,000P가 즉시 지급됩니다."),
    ("반품 고객의 적립은?", "적립된 포인트는 본사(영업마케팅실)로 알려주시면 차감 처리합니다."),
]
for i, (q, a) in enumerate(faqs):
    x = 0.65 + (i % 2) * 6.2
    y = 1.85 + (i // 2) * 1.5
    box(s, x, y, 5.95, 1.3, fill=WHITE, line=LINE)
    text(s, x + 0.3, y + 0.18, 5.4, 0.35, [[("Q. ", {"color": TEAL, "bold": True}), (q, {"bold": True, "color": NAVY})]], size=13.5)
    text(s, x + 0.3, y + 0.58, 5.4, 0.65, a, size=12, color=MUTED, line_spacing=1.25)
box(s, 0.65, 5.05, 12.05, 1.75, fill=RGBColor(0xFB, 0xEC, 0xEC))
text(s, 1.0, 5.28, 11.5, 0.35, "개인정보 — 반드시 지켜주세요", size=14.5, color=RED, bold=True)
text(s, 1.0, 5.7, 11.5, 1.0,
     [["·  매장용 페이지 토큰과 고객 정보(전화번호·이름)는 외부에 공유 금지"],
      ["·  고객 정보는 적립·차감 업무에만 사용하고, 화면을 고객 외 타인에게 보여주지 않기"],
      ["·  고객이 탈퇴·동의 철회를 요청하면 본사(영업마케팅실)로 즉시 전달"]],
     size=13, color=TEXT, line_spacing=1.35)

# ───────────────────────── 8. 마무리
s = slide(NAVY)
text(s, 0.9, 2.1, 11.5, 0.9, "오늘부터 시작하세요!", size=38, color=WHITE, bold=True)
text(s, 0.9, 3.05, 11.5, 0.5, "QR 포스터 비치 → 매장용 페이지 설정 → 첫 고객 가입 안내", size=16, color=RGBColor(0xAF, 0xB9, 0xCE))
links = [
    ("매장용 페이지", "wonsosw-cmd.github.io/eastend-members/staff.html"),
    ("QR 포스터 (인쇄)", "wonsosw-cmd.github.io/eastend-members/qr.html"),
    ("문의", "영업마케팅실  ·  wonsosw@eastend.co.kr"),
]
for i, (k, v) in enumerate(links):
    y = 4.0 + i * 0.85
    box(s, 0.9, y, 11.5, 0.7, fill=RGBColor(0x25, 0x36, 0x5C))
    text(s, 1.25, y + 0.17, 2.6, 0.4, k, size=13.5, color=AMBER, bold=True)
    text(s, 3.9, y + 0.15, 8.3, 0.4, v, size=14.5, color=WHITE)

prs.save(OUT)
print("saved:", OUT)
